# -*- coding: utf-8 -*-
import tkinter as tk
import win32com
# <editor-fold desc="'''클릭 시 반응되는 코드'''">
def getElement(event):
    selection = event.widget.curselection()
    index = selection[0]
    value = event.widget.get(index)

    result.set(value)
    print(index, ' -> ', value)


root = tk.Tk()
root.title("Code4Example.com")  # Add a title

result = tk.StringVar()

tk.Label(root, text="""Click Listbox Element""").grid(row=0, column=0)
tk.Label(root, text="""result""", textvariable=result).grid(row=1, column=0)

var2 = tk.StringVar()
var2.set(('Apple', 'Banana', 'Pear', 'Peach'))
lb = tk.Listbox(root, listvariable=var2)
lb.grid(row=0, column=1)
lb.bind('<<ListboxSelect>>', getElement)  # Select click

root.mainloop()
# </editor-fold>

# <editor-fold desc="검색에 따른 필터링 list박스">
'''검색에 따른 필터링 결과'''
from tkinter import *

# First create application class


class Application(Frame):

    def __init__(self, master=None):
        Frame.__init__(self, master)

        self.pack()
        self.create_widgets()

    # Create main GUI window
    def create_widgets(self):
        self.search_var = StringVar()
        self.search_var.trace("w", self.update_list)
        self.entry = Entry(self, textvariable=self.search_var, width=13)
        self.lbox = Listbox(self, width=45, height=15)

        self.entry.grid(row=0, column=0, padx=10, pady=3)
        self.lbox.grid(row=1, column=0, padx=10, pady=3)

        # Function for updating the list/doing the search.
        # It needs to be called here to populate the listbox.
        self.update_list()

    def update_list(self, *args):
        search_term = self.search_var.get()

        # Just a generic list to populate the listbox
        lbox_list = ['Adam', 'Lucy', 'Barry', 'Bob',
                     'James', 'Frank', 'Susan', 'Amanda', 'Christie']

        self.lbox.delete(0, END)

        for item in lbox_list:
                if search_term.lower() in item.lower():
                    self.lbox.insert(END, item)


root = Tk()
root.title('Filter Listbox Test')
app = Application(master=root)
print ('Starting mainloop()')
app.mainloop()
# </editor-fold>

# <editor-fold desc="이미지삽입">
'''이미지삽입 - 필로우 사용해서'''
#Import tkinter library
from tkinter import *
from PIL import Image,ImageTk
#Create an instance of tkinter frame
win = Tk()
#Set the geometry
win.geometry("750x550")
#Load the image
img= Image.open("tutorialspoint.jpg")
#Convert To photoimage
tkimage= ImageTk.PhotoImage(img)
#Display the Image
label=Label(win,image=tkimage)
label.pack()
win.mainloop()
# </editor-fold>

# <editor-fold desc="슬라이드 복사 붙여넣기">
'''pptx-슬라이드 복사 붙여넣기'''
from pptx import Presentation
import copy
import pythoncom

'''슬라이드 복사 함수'''


def duplicate_slide(prs, index):
    # 원본 슬라이드 번호 정의

    source_slide = prs.slides[index]
    # 빈 레이아웃 설정 (슬라이드 레이아웃 정의 필요)
    try:
        slide_layout = prs.slide_layouts[6]
    except:
        slide_layout = prs.slide_layouts[len(prs.slide_layouts)]
    copied_slide = prs.slides.add_slide(slide_layout)
    # 원본 슬라이드에서 shape 복제
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    return copied_slide
# </editor-fold>

# <editor-fold desc="pptx 텍스트 바꾸기 1">
'''pptx 텍스트 바꾸기 1'''
prs = Presentation('blah.pptx')

# To get shapes in your slides
slides = [slide for slide in prs.slides]
shapes = []
for slide in slides:
    for shape in slide.shapes:
        shapes.append(shape)

def replace_text(self, replacements: dict, shapes: List):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text = new_text

replace_text({'string to replace': 'replacement text'}, shapes)
# </editor-fold>

# <editor-fold desc="pptx텍스트 바꾸기 2">
'''pptx 텍스트 바꾸기 2'''
def search_and_replace(search_str, repl_str, input, output):
    """"search and replace text in PowerPoint while preserving formatting"""
    #Useful Links ;)
    #https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    #https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
    from pptx import Presentation
    prs = Presentation(input)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(search_str))!=-1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text
    prs.save(output)
# </editor-fold>

# <editor-fold desc="pptx텍스트 바꾸기 3">
'''pptx 텍스트 바꾸기 3'''
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
prs = Presentation("./template.pptx")
'''모든 슬라이드에서 텍스트 찾아 바꾸기'''
for slide in prs.slides:
    for shape in slide.shapes:
        # 텍스트박스에서
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if 찾을내용 in run.text:
                        run.text = run.text.replace(찾을내용, str(바꿀내용))
        # 그룹된 텍스트박스에서
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shp in shape.shapes:
                if shp.has_text_frame:
                    for paragraph in shp.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if 찾을내용 in run.text:
                                run.text = run.text.replace(찾을내용, str(바꿀내용))
        # 표에서
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if 찾을내용 in run.text:
                                run.text = run.text.replace(찾을내용, str(바꿀내용))

prs.save('./output.pptx')
# </editor-fold>

# <editor-fold desc="텍스트박스 생성 및 위치 설정">
'''텍스트 박스 생성 및 위치 설정'''
from pptx import Presentation
from pptx.util import Inches, Pt,Cm
from pptx.enum.text import PP_ALIGN


prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = Cm(3)
top = Cm(2.5)
width = Cm(15)
height = Cm(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Hello"
txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Just an example"
font = run.font


prs.save('test.pptx')
# </editor-fold>

# <editor-fold desc="폰트설정">
'''폰트설정 - 테이블 제외'''
from google.cloud import translate_v2 as translate
import os
from pptx import Presentation

os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = r"Google-Cloud Key가 저장된 파일 경로"
client = translate.Client()

prs=Presentation("Trans_Sample.pptx")
Slides = prs.slides
for slide in Slides:
    Shapes = slide.shapes
    for shape in Shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                font=paragraph.runs[0].font
                Fname = font.name
                Fsize = font.size
                cur_text = paragraph.text
                result = client.translate(cur_text, target_language='en')
                new_text = result['translatedText']
                paragraph.text = new_text
                for run in paragraph.runs:
                    font = run.font
                    font.name = Fname
                    font.size = Fsize
prs.save('T_Result.pptx')
# </editor-fold>

# <editor-fold desc="폰트설정 - 테이블">
'''폰트설정 테이블 포함'''
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def C_Font (shp,Fname):
    for paragraph in shp.text_frame.paragraphs:
         for run in paragraph.runs:
            run.font.name = Fname

def AutoFont(file,Tpath,Rpath,Fname):
    prs = Presentation(Tpath + '/' + file)
    for slide in prs.slides:
        for shape in slide.shapes:
             if shape.has_text_frame:
                 C_Font(shape,Fname)
             if shape.has_table:
                 for row in shape.table.rows:
                     for cell in row.cells:
                         C_Font(cell,Fname)
             if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                 for shp in shape.shapes:
                     if shp.has_text_frame:
                         C_Font(shp,Fname)
    prs.save(Rpath + '/' + file)

Tpath = 'Sample'
Rpath = 'Result'
Fname = '맑은 고딕'
file = 'Sample.pptx'

AutoFont(file,Tpath,Rpath,Fname)
# </editor-fold>

# <editor-fold desc="파워포인트 이미지로 추출">
file_address="C:\\pptx_테스트_파일들\\pptx2\\2022_KOREG_IOKI_전략_제안서_v0.1.pptx"
def PPTX_image_export(file_address,file_name,pages):
    import win32com.client, sys
    '''
    orignial_address:원래경로
    file_address:경로
    file_name:경로 내 파일이름
    pages:해당하는 페이지(list나 range형식으로 제공해야함)
    '''
    orignial_address=os.getcwd()
    os.chdir(file_address)
    FILENAME = file_name

    pythoncom.CoInitialize()
    APPLICATION = win32com.client.Dispatch("PowerPoint.Application")
    PRESENTATION = APPLICATION.Presentations.Open(FILENAME, ReadOnly=False)  # 수정필요

    for page, slide in enumerate(PRESENTATION.Slides):
        print(page)
        try:
            if page in pages:
                PRESENTATION.Slides[page].Export(r"C:\\pptx_테스트_파일들\\{0}.jpg".format(page), "JPG")
            else:
                continue
        except pythoncom.error:
            print("error")

    APPLICATION.Quit()

    PRESENTATION =  None
    APPLICATION = None
    pythoncom.CoUninitialize()

    os.chdir(orignial_address)
# </editor-fold>

# <editor-fold desc="메뉴만들기">
import tkinter
window=tkinter.Tk()
window.title("YUN DAE HEE")
window.geometry("640x480+100+100")
window.resizable(False, False)

def close():
    window.quit()
    window.destroy()

menubar=tkinter.Menu(window)

menu_1=tkinter.Menu(menubar, tearoff=0)
menu_1.add_command(label="하위 메뉴 1-1")
menu_1.add_command(label="하위 메뉴 1-2")
menu_1.add_separator()
menu_1.add_command(label="하위 메뉴 1-3", command=close)
menubar.add_cascade(label="상위 메뉴 1", menu=menu_1)

menu_2=tkinter.Menu(menubar, tearoff=0, selectcolor="red")
menu_2.add_radiobutton(label="하위 메뉴 2-1", state="disable")
menu_2.add_radiobutton(label="하위 메뉴 2-2")
menu_2.add_radiobutton(label="하위 메뉴 2-3")
menubar.add_cascade(label="상위 메뉴 2", menu=menu_2)

menu_3=tkinter.Menu(menubar, tearoff=0)
menu_3.add_checkbutton(label="하위 메뉴 3-1")
menu_3.add_checkbutton(label="하위 메뉴 3-2")
menubar.add_cascade(label="상위 메뉴 3", menu=menu_3)

window.config(menu=menubar)

window.mainloop()
# </editor-fold>

'''슬라이드 합치는건 기본적으로 거기 있는 shaeps들 전부 복사하여 붙여넣는 방식인듯.'''